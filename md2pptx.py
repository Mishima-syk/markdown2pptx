from markdown import Markdown
from html.parser import HTMLParser
from pptx import Presentation
from pptx.util import Inches, Pt
import pptx
import click
from logging import getLogger, StreamHandler, DEBUG
logger = getLogger(__name__)
handler = StreamHandler()
handler.setLevel(DEBUG)
logger.setLevel(DEBUG)
logger.addHandler(handler)
logger.propagate = False


def class2dic(attr):
    d = {}
    for k, v in attr:
        if k == "class":
            for cls in v.split():
                ek, ev = cls.split("-")
                d[ek] = float(ev)
        else:
            d[k] = v
    return d


class MyHTMLParser(HTMLParser):
    def __init__(self):
        super(MyHTMLParser, self).__init__()
        self.prs = Presentation()
        self.tags = []
        self.focus = None    # current item
        self.ln = 0          # layout number
        self.header = False  # TF BF marker
        self.slide = None    # current slide
        self.pg = None       # current paragraph

    def handle_starttag(self, tag, attrs):
        attr_dict = class2dic(attrs)
        if tag == "h1" or tag == "h2":
            hln = int(tag[-1]) - 1
            ln = int(attr_dict.get("layout", hln))
            self.ln = ln
            slide = self.prs.slides.add_slide(self.prs.slide_layouts[ln])
            self.slide = slide
            self.focus = slide
            self.tags.append(tag)
        else:
            self.tags.append(tag)
        getattr(self, "handle_starttag_layout{}".format(self.ln))(tag, attr_dict)

    def handle_endtag(self, tag):
        getattr(self, "handle_endtag_layout{}".format(self.ln))(tag)

    def handle_data(self, data):
        getattr(self, "handle_data_layout{}".format(self.ln))(data)

    # layout 0 #

    def handle_starttag_layout0(self, tag, attrs):
            pass

    def handle_endtag_layout0(self, tag):
        pass

    def handle_data_layout0(self, data):
        tag = self.tags.pop()
        if tag == "h1":
                self.focus.shapes.title.text = data
        elif tag == "p":
            self.focus.placeholders[1].text = data
        else:
            print("Not Implemented...")

    # layout 1 #

    def handle_starttag_layout1(self, tag, attrs):
        if tag == "h2":
            pass
        else:
            self.tags.append(tag)

    def handle_endtag_layout1(self, tag):
        pass

    def handle_data_layout1(self, data):
        tag = self.tags.pop()
        if tag == "h2":
            self.focus.shapes.title.text = data
        elif tag == "h3":
            p = self.focus.shapes.placeholders[1].text_frame.add_paragraph()
            p.text = data
        elif tag == "h4":
            p = self.focus.shapes.placeholders[1].text_frame.add_paragraph()
            p.text = data
            p.level = 1
        elif tag == "h5":
            p = self.focus.shapes.placeholders[1].text_frame.add_paragraph()
            p.text = data
            p.level = 2
        elif tag == "h6":
            p = self.focus.shapes.placeholders[1].text_frame.add_paragraph()
            p.text = data
            p.level = 3
        else:
            print("Not Implemented...")

    # layout 2 #

    def handle_starttag_layout2(self, tag, attr_dict):
            print("Not Implemented...")

    def handle_endtag_layout2(self, tag):
        pass

    def handle_data_layout2(self, data):
        tag = self.tags.pop()
        print("Not Implemented...")

    # layout 3 #

    def handle_starttag_layout3(self, tag, attr_dict):
            print("Not Implemented...")

    def handle_endtag_layout3(self, tag):
        pass

    def handle_data_layout3(self, data):
        tag = self.tags.pop()
        print("Not Implemented...")

    # layout 4 #

    def handle_starttag_layout4(self, tag, attr_sict):
            print("Not Implemented...")

    def handle_endtag_layout4(self, tag):
        pass

    def handle_data_layout4(self, data):
        tag = self.tags.pop()
        print("Not Implemented...")

    # layout 5 #

    def handle_starttag_layout5(self, tag, attr_dict):
            print("Not Implemented...")

    def handle_endtag_layout5(self, tag):
        pass

    def handle_data_layout5(self, data):
        tag = self.tags.pop()
        print("Not Implemented...")

    # layout 6 #

    def handle_starttag_layout6(self, tag, attr_dict):
        if tag == "img":
            print(attr_dict)
            img_path = attr_dict.get("src", None)
            left = Inches(attr_dict.get("left", 1))
            top = Inches(attr_dict.get("top", 1))
            height = attr_dict.get("height", None)
            width = attr_dict.get("width", None)
            if height is None and width is None:
                self.slide.shapes.add_picture(img_path, left, top)
            elif height is None:
                self.slide.shapes.add_picture(img_path, left, top, width=Inches(width))
            elif width is None:
                self.slide.shapes.add_picture(img_path, left, top, height=Inches(height))
            else:
                self.slide.shapes.add_picture(img_path, left, top, height=Inches(height), width=Inches(width))
        elif tag == "p":
            left = attr_dict.get("left", None)
            top = attr_dict.get("top", None)
            right = attr_dict.get("right", None)
            bottom = attr_dict.get("bottom", None)
            fontsize = attr_dict.get("fontsize", None)
            if left is None and top is None and right is None and bottom is None:
                if isinstance(self.focus, pptx.text.text.TextFrame):
                    self.pg = self.focus.add_paragraph()
                    if fontsize is not None:
                        self.pg.font.size = Pt(int(fontsize))
                self.tags.append(tag)
            else:
                tb = self.slide.shapes.add_textbox(Inches(left),
                                                   Inches(top),
                                                   Inches(right),
                                                   Inches(bottom))
                self.focus = tb.text_frame
                self.header = True
                self.tags.append(tag)
        else:
            self.tags.append(tag)

    def handle_data_layout6(self, data):
        tag = self.tags.pop()
        if tag == "h2":
            pass
        elif tag == "p":
            if isinstance(self.focus, pptx.text.text.TextFrame):
                if self.header:
                    self.focus.text = data
                else:
                    self.pg.text = data
        elif tag == "strong":
            p = self.focus.add_paragraph()
            p.text = data
            p.font.bold = True
        else:
            print("Not Implemented...")

    def handle_endtag_layout6(self, tag):
        if tag == "p":
            if self.header:
                self.header = False

    # layout 7 #

    def handle_starttag_layout7(self, tag, attr_dict):
            print("Not Implemented...")

    def handle_endtag_layout7(self, tag):
        pass

    def handle_data_layout7(self, data):
        tag = self.tags.pop()
        print("Not Implemented...")

    # layout 8 #

    def handle_starttag_layout8(self, tag, attr_dict):
            print("Not Implemented...")

    def handle_endtag_layout8(self, tag):
        pass

    def handle_data_layout8(self, data):
        tag = self.tags.pop()
        print("Not Implemented...")

    # layout 9 #

    def handle_starttag_layout9(self, tag, attr_dict):
            print("Not Implemented...")

    def handle_endtag_layout9(self, tag):
        pass

    def handle_data_layout9(self, data):
        tag = self.tags.pop()
        print("Not Implemented...")

    def close(self, output="sample.pptx"):
        self.prs.save(output)


@click.command()
@click.argument('input')
@click.option('--output', '-o', default='output.pptx', help='Name of output_file.')
def cli(input, output):
    md = Markdown(extensions=['markdown.extensions.attr_list'])
    parser = MyHTMLParser()
    md_txt = open(input, "r").read()
    logger.debug(md.convert(md_txt))
    parser.feed(md.convert(md_txt).replace("\n", ""))  # ???
    parser.close(output)


if __name__ == '__main__':
    cli()
